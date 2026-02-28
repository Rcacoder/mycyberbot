from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()

    # Define some theme colors (Dark mode aesthetic from web UI)
    dark_bg = RGBColor(13, 15, 23)
    accent_blue = RGBColor(0, 242, 254)
    text_main = RGBColor(226, 232, 240)
    
    # --- Slide 1: Title Slide ---
    # Layout 0 is usually the title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Daily Cyber Attack Reporter"
    title.text_frame.paragraphs[0].font.color.rgb = accent_blue
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle.text = "Automated Threat Intelligence & Academic Integration\nUsing AI & RSS Aggregation"
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_main
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # --- Slide 2: Architecture & Features ---
    # Layout 1 is usually Title and Content
    slide_layout_2 = prs.slide_layouts[1]
    slide_2 = prs.slides.add_slide(slide_layout_2)
    
    # Set background color
    background_2 = slide_2.background
    fill_2 = background_2.fill
    fill_2.solid()
    fill_2.fore_color.rgb = dark_bg
    
    title_2 = slide_2.shapes.title
    title_2.text = "System Architecture & Features"
    title_2.text_frame.paragraphs[0].font.color.rgb = accent_blue
    title_2.text_frame.paragraphs[0].font.bold = True
    
    body_shape = slide_2.placeholders[1]
    tf = body_shape.text_frame
    tf.clear() # clear existing paragraphs
    
    # Add bullet points
    bullets = [
        ("📡 Data Aggregation", "Continuously fetches the latest cybersecurity news from top RSS feeds (The Hacker News, CISA, etc.)."),
        ("🧠 DeepSeek AI Analysis", "Identifies the Top 10 global threats and generates university-level teaching lessons for the Top 2 incidents."),
        ("✉️ Automated Notification", "Formats insights into a clean HTML email delivered daily via SMTP."),
        ("🖥️ Interactive Web Dashboard", "Stores historical JSON reports accessible via a premium glassmorphism FastAPI web interface.")
    ]
    
    for i, (heading, desc) in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = f"{heading}: "
        p.font.bold = True
        p.font.color.rgb = accent_blue
        p.font.size = Pt(20)
        p.level = 0
        
        # Add description text in the same paragraph
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = text_main
        run.font.size = Pt(20)

    # Save presentation
    filepath = "Cyber_Attack_Reporter_Presentation.pptx"
    prs.save(filepath)
    print(f"Successfully generated {filepath}")

if __name__ == "__main__":
    create_presentation()
